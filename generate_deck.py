"""
SBI Dost - Hackathon Idea Deck Generator
Generates a professional 10-slide PowerPoint presentation for the SBI Hackathon @ GFF 2026.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==========================================
# COLOR PALETTE
# ==========================================
NAVY      = RGBColor(15, 23, 42)       # Deep Slate Navy
DARK_BLUE = RGBColor(30, 41, 59)       # Slate 800
SBI_BLUE  = RGBColor(0, 102, 204)      # Corporate SBI Blue
LIGHT_BLUE= RGBColor(56, 189, 248)     # Accent Cyan
WHITE     = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(248, 250, 252)     # Slate 50
CREAM     = RGBColor(252, 246, 242)     # Peach Cream
CHARCOAL  = RGBColor(43, 34, 31)       # Warm Charcoal
WARM_GRAY = RGBColor(123, 108, 101)    # Warm Secondary Text
SUCCESS   = RGBColor(16, 185, 129)     # Emerald
DANGER    = RGBColor(239, 68, 68)      # Rose Red
PURPLE    = RGBColor(139, 92, 246)     # Violet
AMBER     = RGBColor(245, 158, 11)     # Amber

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ==========================================
# HELPER FUNCTIONS
# ==========================================

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    """Add a rectangle shape with optional fill and border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    line = shape.line
    if border_color:
        line.color.rgb = border_color
        line.width = Pt(border_width or 1)
    else:
        line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=CHARCOAL, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=CHARCOAL, line_spacing=1.4, bold_prefix=True):
    """Add a bulleted list text box. Items can be strings or (bold_part, rest) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        
        if isinstance(item, tuple):
            # Bold prefix + normal text
            run1 = p.add_run()
            run1.text = "• " + item[0]
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = 'Calibri'
            
            run2 = p.add_run()
            run2.text = " " + item[1]
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = WARM_GRAY
            run2.font.name = 'Calibri'
        else:
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = 'Calibri'
    return txBox

def dark_background(slide):
    """Fill the slide background with deep navy."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

def cream_background(slide):
    """Fill the slide background with warm peach cream."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = CREAM

def add_footer_bar(slide):
    """Add a consistent dark footer bar."""
    bar = add_shape(slide, Inches(0), SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55), fill_color=DARK_BLUE)
    add_text_box(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(5), Inches(0.35),
                 "SBI Dost  •  SBI Hackathon @ GFF 2026  •  Pillar: Digital Engagement",
                 font_size=9, color=RGBColor(148, 163, 184), font_name='Calibri')
    add_text_box(slide, Inches(9), SLIDE_H - Inches(0.45), Inches(4), Inches(0.35),
                 "Theme: Agentic AI & Emerging Tech",
                 font_size=9, color=RGBColor(148, 163, 184), alignment=PP_ALIGN.RIGHT, font_name='Calibri')

def add_section_badge(slide, left, top, text, bg_color=SBI_BLUE):
    """Add a small colored badge label."""
    shape = add_shape(slide, left, top, Inches(2.2), Inches(0.32), fill_color=bg_color)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'

def add_metric_card(slide, left, top, width, height, label, value, sub_text="", accent_color=SBI_BLUE):
    """Add a metric card with label, big number, and optional subtitle."""
    card = add_shape(slide, left, top, width, height, fill_color=WHITE, border_color=RGBColor(235, 220, 211), border_width=1)
    # Accent bar at top
    add_shape(slide, left, top, width, Inches(0.06), fill_color=accent_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.3),
                 label, font_size=9, color=WARM_GRAY, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), Inches(0.5),
                 value, font_size=24, color=CHARCOAL, bold=True)
    if sub_text:
        add_text_box(slide, left + Inches(0.2), top + Inches(1.0), width - Inches(0.4), Inches(0.3),
                     sub_text, font_size=9, color=WARM_GRAY)

# ==========================================
# SLIDE 1: COVER
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
dark_background(slide)

# Large decorative circle accent (top right)
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(30, 41, 59)
circle.line.fill.background()

# Title
add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(1),
             "SBI DOST", font_size=52, bold=True, color=WHITE)

# Subtitle
add_text_box(slide, Inches(1), Inches(2.6), Inches(8), Inches(0.6),
             "Autonomous Personal Finance Concierge & Agentic Advisor",
             font_size=22, color=LIGHT_BLUE)

# Tagline
add_text_box(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.8),
             "Unlocking Digital Engagement through Proactive Agentic Banking",
             font_size=16, color=RGBColor(148, 163, 184))

# Separator line
add_shape(slide, Inches(1), Inches(4.5), Inches(2), Inches(0.04), fill_color=SBI_BLUE)

# Event info
add_text_box(slide, Inches(1), Inches(4.8), Inches(6), Inches(0.4),
             "SBI Hackathon @ Global Fintech Fest 2026",
             font_size=14, color=RGBColor(148, 163, 184))
add_text_box(slide, Inches(1), Inches(5.3), Inches(6), Inches(0.4),
             "Pillar: Digital Engagement  •  Theme: Agentic AI & Emerging Tech",
             font_size=12, color=RGBColor(100, 116, 139))

# Team info
add_text_box(slide, Inches(1), Inches(6.1), Inches(6), Inches(0.4),
             "Team: Jatin Garg", font_size=13, color=WHITE, bold=True)

add_footer_bar(slide)

# ==========================================
# SLIDE 2: THE PROBLEM
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
cream_background(slide)

add_section_badge(slide, Inches(0.8), Inches(0.5), "THE PROBLEM", DANGER)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "The Digital Engagement Gap in Indian Banking",
             font_size=30, bold=True, color=CHARCOAL)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(10), Inches(0.5),
             "Banks are losing customers, revenue, and trust due to passive, reactive digital experiences.",
             font_size=14, color=WARM_GRAY)

# Problem Cards (4 cards in a row)
problems = [
    ("Silent Financial Leakage", "Customers lose ₹10,000+ annually to forgotten recurring mandates (streaming, gym, cloud storage) that auto-debit without usage.", "💸", DANGER),
    ("Passive Account Engagement", "Idle cash sits in savings accounts earning 3.0% when SBI FDs/MFs offer 6-8%. Banks miss cross-sell due to lack of proactive outreach.", "📉", AMBER),
    ("Reactive Budget Alerts", "Standard apps notify AFTER overspend. No mechanism exists to predict and prevent budget breaches in real-time.", "⚠️", PURPLE),
    ("Support Desk Overload", "Multi-step queries (mandate blocking, audit, transfers) overwhelm call centers. 60% of tickets are automatable.", "📞", SBI_BLUE),
]

for i, (title, desc, emoji, color) in enumerate(problems):
    x = Inches(0.8 + i * 3.05)
    y = Inches(2.6)
    card = add_shape(slide, x, y, Inches(2.85), Inches(4.0), fill_color=WHITE, border_color=RGBColor(235, 220, 211), border_width=1)
    # Top accent
    add_shape(slide, x, y, Inches(2.85), Inches(0.06), fill_color=color)
    # Emoji
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.6), Inches(0.6),
                 emoji, font_size=28)
    # Title
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.45), Inches(0.5),
                 title, font_size=15, bold=True, color=CHARCOAL)
    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(1.5), Inches(2.45), Inches(2.2),
                 desc, font_size=11, color=WARM_GRAY)

add_footer_bar(slide)

# ==========================================
# SLIDE 3: THE SOLUTION
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
cream_background(slide)

add_section_badge(slide, Inches(0.8), Inches(0.5), "THE SOLUTION", SUCCESS)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "SBI Dost: Proactive, Secure, Agentic Banking",
             font_size=30, bold=True, color=CHARCOAL)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(10), Inches(0.5),
             "An autonomous AI concierge that monitors, advises, and acts — with full user control and transparency.",
             font_size=14, color=WARM_GRAY)

features = [
    ("🔍 Proactive Mandate Auditor", 
     "Continuously monitors recurring auto-debit mandates. Detects subscriptions unused for 60/90/120 days. Recommends and executes cancellations under RBI guidelines (Policy S-812).",
     SUCCESS),
    ("🧠 RAG-Powered Policy Engine", 
     "Before executing any transfer, the agent queries a vector database of SBI policies. Validates transfer limits (Policy L-204: ₹20K daily cap), OTP thresholds, and beneficiary compliance.",
     SBI_BLUE),
    ("🛡️ Human-in-the-Loop Guardrails", 
     "All monetary actions require explicit user authorization via a visual OTP approval card. The agent plans, verifies, then HALTS for human confirmation before executing core banking APIs.",
     DANGER),
    ("📊 Reasoning Transparency Console", 
     "Displays the agent's THOUGHT → ACTION → OBSERVATION chain directly to the user. Builds algorithmic trust by showing exactly how decisions are made — no black boxes.",
     PURPLE),
]

for i, (title, desc, color) in enumerate(features):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(2.5 + row * 2.2)
    
    card = add_shape(slide, x, y, Inches(5.7), Inches(2.0), fill_color=WHITE, border_color=RGBColor(235, 220, 211), border_width=1)
    # Left accent strip
    add_shape(slide, x, y, Inches(0.06), Inches(2.0), fill_color=color)
    
    add_text_box(slide, x + Inches(0.25), y + Inches(0.2), Inches(5.2), Inches(0.4),
                 title, font_size=16, bold=True, color=CHARCOAL)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.7), Inches(5.2), Inches(1.2),
                 desc, font_size=12, color=WARM_GRAY)

add_footer_bar(slide)

# ==========================================
# SLIDE 4: HOW IT WORKS (ReAct Loop)
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
cream_background(slide)

add_section_badge(slide, Inches(0.8), Inches(0.5), "HOW IT WORKS", SBI_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "The ReAct (Reasoning & Action) Agentic Loop",
             font_size=30, bold=True, color=CHARCOAL)

# Flow steps as connected cards
steps = [
    ("1", "User Query\nor Trigger", "Natural language input\nor proactive anomaly", SBI_BLUE),
    ("2", "Policy RAG\nLookup", "search_sbi_policy()\nfetches banking limits", PURPLE),
    ("3", "Balance &\nRecipient Check", "get_balance()\ncheck_recipient()", AMBER),
    ("4", "Reasoning\nTrace Log", "THOUGHT → ACTION\n→ OBSERVATION", RGBColor(100, 116, 139)),
    ("5", "HITL Security\nGuardrail", "OTP Approval\nRequired > ₹10,000", DANGER),
    ("6", "API Execution\n& Update", "execute_transfer()\nBudgets refreshed", SUCCESS),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.08)
    y = Inches(2.2)
    
    # Card
    card = add_shape(slide, x, y, Inches(1.88), Inches(2.8), fill_color=WHITE, border_color=RGBColor(235, 220, 211), border_width=1)
    add_shape(slide, x, y, Inches(1.88), Inches(0.06), fill_color=color)
    
    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.64), y + Inches(0.2), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    # Title
    add_text_box(slide, x + Inches(0.1), y + Inches(1.0), Inches(1.68), Inches(0.7),
                 title, font_size=13, bold=True, color=CHARCOAL, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x + Inches(0.1), y + Inches(1.7), Inches(1.68), Inches(0.9),
                 desc, font_size=10, color=WARM_GRAY, alignment=PP_ALIGN.CENTER, font_name='Consolas')

    # Arrow between cards
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(1.88), y + Inches(1.1), Inches(0.25), Inches(0.4),
                     "→", font_size=20, bold=True, color=RGBColor(148, 163, 184), alignment=PP_ALIGN.CENTER)

# Example workflow box
add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.4), fill_color=DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(5.4), Inches(2), Inches(0.3),
             "LIVE EXAMPLE", font_size=10, bold=True, color=LIGHT_BLUE)
add_text_box(slide, Inches(1.0), Inches(5.8), Inches(11.2), Inches(0.8),
             'User: "Settle my rent of ₹12,000"  →  THOUGHT: Check balance  →  ACTION: get_balance() → 45,000  →  '
             'THOUGHT: Within L-204 limits but > 10K threshold  →  ACTION: search_sbi_policy()  →  '
             'HITL: OTP Card Shown  →  User Approves  →  ACTION: execute_transfer() → Success  →  Balance: ₹33,000',
             font_size=11, color=RGBColor(203, 213, 225), font_name='Consolas')

add_footer_bar(slide)

# ==========================================
# SLIDE 5: TECHNOLOGY STACK
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
cream_background(slide)

add_section_badge(slide, Inches(0.8), Inches(0.5), "TECHNOLOGY STACK", PURPLE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Architecture & Technology Components",
             font_size=30, bold=True, color=CHARCOAL)

tech_items = [
    ("Frontend UI", "Vite + React\nVanilla CSS", "Tabbed multi-page SaaS dashboard\nwith hybrid dark-peach theme,\nglassmorphic panels, micro-animations", SBI_BLUE),
    ("Agent Core Engine", "ReAct Loop\n(JavaScript)", "Custom Reasoning & Action loop:\nThought → Plan → Tool Call\n→ Result → Response", SUCCESS),
    ("Policy RAG", "Vector Database\n(Chroma / pgvector)", "Indexes SBI banking policies,\nFAQs, transfer circulars for\ndynamic policy verification", PURPLE),
    ("Security Layer", "HITL Guardrails\n+ OTP Validation", "Blocks all debit APIs until\nexplicit user authorization.\nPrompt injection guards.", DANGER),
    ("Mock APIs", "SBI Core Banking\nSimulator", "Accounts, Transactions,\nBudgets, Mandates, Beneficiaries\n— fully stateful sandbox", AMBER),
]

for i, (title, tech, desc, color) in enumerate(tech_items):
    y_pos = Inches(2.0 + i * 1.02)
    
    # Row bar
    card = add_shape(slide, Inches(0.8), y_pos, Inches(11.7), Inches(0.9), fill_color=WHITE, border_color=RGBColor(235, 220, 211), border_width=1)
    add_shape(slide, Inches(0.8), y_pos, Inches(0.06), Inches(0.9), fill_color=color)
    
    add_text_box(slide, Inches(1.1), y_pos + Inches(0.15), Inches(2.2), Inches(0.6),
                 title, font_size=14, bold=True, color=CHARCOAL)
    add_text_box(slide, Inches(3.5), y_pos + Inches(0.15), Inches(2.5), Inches(0.6),
                 tech, font_size=11, color=SBI_BLUE, bold=True, font_name='Consolas')
    add_text_box(slide, Inches(6.5), y_pos + Inches(0.05), Inches(5.8), Inches(0.8),
                 desc, font_size=11, color=WARM_GRAY)

add_footer_bar(slide)

# ==========================================
# SLIDE 6: USE CASE 1 - MANDATE LEAKAGE
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
cream_background(slide)

add_section_badge(slide, Inches(0.8), Inches(0.5), "USE CASE 1", DANGER)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Proactive Subscription Leakage Detection & Mandate Blocking",
             font_size=26, bold=True, color=CHARCOAL)

# Left column: Flow
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
             "Agent Workflow (ReAct Trace):", font_size=14, bold=True, color=CHARCOAL)

workflow_steps = [
    ("GOAL:", "Scan active recurring mandates for utilization leakage"),
    ("THOUGHT:", "Checking RBI guidelines for mandate cancellation rights"),
    ("ACTION:", "search_sbi_policy('mandate cancellation') → Policy S-812: Customer can block mandates instantly"),
    ("ACTION:", "get_subscriptions() → Found 3 active mandates"),
    ("THOUGHT:", "Streaming Premium unused for 120 days (₹699/mo). Gym unused for 75 days (₹1,500/mo)"),
    ("RESPONSE:", "Recommend blocking Streaming Premium. Annual savings: ₹8,388"),
    ("ACTION:", "cancel_subscription('Streaming Premium') → Success"),
]

for i, (tag, text) in enumerate(workflow_steps):
    y = Inches(2.5 + i * 0.55)
    
    tag_color = SBI_BLUE
    if "THOUGHT" in tag: tag_color = SUCCESS
    elif "ACTION" in tag: tag_color = PURPLE
    elif "RESPONSE" in tag: tag_color = AMBER
    elif "GOAL" in tag: tag_color = SBI_BLUE
    
    # Tag
    add_text_box(slide, Inches(0.8), y, Inches(1.2), Inches(0.35),
                 tag, font_size=9, bold=True, color=tag_color, font_name='Consolas')
    # Text
    add_text_box(slide, Inches(2.0), y, Inches(4.5), Inches(0.45),
                 text, font_size=10, color=CHARCOAL, font_name='Consolas')

# Right column: Impact metrics
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.7), fill_color=WHITE, border_color=RGBColor(235, 220, 211), border_width=1)
add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(0.4),
             "Business Impact", font_size=16, bold=True, color=CHARCOAL)

impact_items = [
    ("₹8,388", "Annual savings per customer from blocking 1 leaky mandate"),
    ("120 days", "Average idle time before Dost detects and flags mandate leakage"),
    ("Policy S-812", "RBI mandate allows instant digital cancellation — no vendor clearance needed"),
    ("40%", "Estimated reduction in mandate-related support tickets"),
]

for i, (metric, desc) in enumerate(impact_items):
    y = Inches(2.8 + i * 0.9)
    add_text_box(slide, Inches(7.3), y, Inches(1.8), Inches(0.4),
                 metric, font_size=22, bold=True, color=SBI_BLUE)
    add_text_box(slide, Inches(7.3), y + Inches(0.4), Inches(4.8), Inches(0.4),
                 desc, font_size=11, color=WARM_GRAY)

add_footer_bar(slide)

# ==========================================
# SLIDE 7: USE CASE 2 - RENT TRANSFER
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
cream_background(slide)

add_section_badge(slide, Inches(0.8), Inches(0.5), "USE CASE 2", SBI_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Conversational Fund Transfer with RAG Policy Limits & HITL Security",
             font_size=26, bold=True, color=CHARCOAL)

# Left: The flow
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
             "Agent Workflow (ReAct Trace):", font_size=14, bold=True, color=CHARCOAL)

rent_steps = [
    ("GOAL:", "Transfer ₹12,000 for monthly rent payment"),
    ("ACTION:", "get_balance() → Balance: ₹45,000 ✓"),
    ("ACTION:", "search_sbi_policy('transfer limits') → Policy L-204: Daily cap ₹20,000. OTP required > ₹10,000"),
    ("THOUGHT:", "₹12,000 is within daily cap but exceeds ₹10K OTP threshold"),
    ("ACTION:", "check_recipient('Landlord') → Acc: XXXXXX3210, SBI, IFSC: SBIN0000001 ✓"),
    ("⚠️ HITL:", "Authorization prompt dispatched → OTP Card shown to user"),
    ("ACTION:", "execute_transfer('Landlord', 12000) → TXN ID: SBI-482916354 ✓"),
    ("RESULT:", "New Balance: ₹33,000. Rent budget category updated."),
]

for i, (tag, text) in enumerate(rent_steps):
    y = Inches(2.5 + i * 0.52)
    
    tag_color = SBI_BLUE
    if "THOUGHT" in tag: tag_color = SUCCESS
    elif "ACTION" in tag: tag_color = PURPLE
    elif "HITL" in tag: tag_color = DANGER
    elif "RESULT" in tag: tag_color = SUCCESS
    
    add_text_box(slide, Inches(0.8), y, Inches(1.2), Inches(0.35),
                 tag, font_size=9, bold=True, color=tag_color, font_name='Consolas')
    add_text_box(slide, Inches(2.0), y, Inches(4.5), Inches(0.45),
                 text, font_size=10, color=CHARCOAL, font_name='Consolas')

# Right: HITL Security Card visual
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.7), fill_color=DARK_BLUE)
add_text_box(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(0.4),
             "🛡️ Human-in-the-Loop Security Guardrail", font_size=15, bold=True, color=WHITE)
add_text_box(slide, Inches(7.3), Inches(2.8), Inches(5), Inches(0.3),
             "Transaction exceeds ₹10,000 threshold (Policy L-204). Explicit authorization required:",
             font_size=11, color=RGBColor(148, 163, 184))

# Mock OTP card
otp_card = add_shape(slide, Inches(7.5), Inches(3.4), Inches(4.8), Inches(2.8), fill_color=RGBColor(30, 41, 59), border_color=RGBColor(239, 68, 68), border_width=2)

labels = [
    ("Action:", "Transfer Funds"),
    ("Recipient:", "Landlord Rent Account"),
    ("Account:", "XXXXXX3210 (SBI)"),
    ("Amount:", "₹12,000"),
]
for i, (label, value) in enumerate(labels):
    y = Inches(3.6 + i * 0.45)
    add_text_box(slide, Inches(7.8), y, Inches(1.2), Inches(0.3),
                 label, font_size=10, color=RGBColor(148, 163, 184), font_name='Consolas')
    add_text_box(slide, Inches(9.2), y, Inches(2.8), Inches(0.3),
                 value, font_size=11, color=WHITE, bold=True, font_name='Consolas')

# Approve button visual
approve = add_shape(slide, Inches(7.8), Inches(5.5), Inches(2.2), Inches(0.45), fill_color=SUCCESS)
tf = approve.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = tf.paragraphs[0].add_run()
run.text = "✓ Authorize & Transfer"
run.font.size = Pt(11)
run.font.bold = True
run.font.color.rgb = WHITE

decline = add_shape(slide, Inches(10.2), Inches(5.5), Inches(1.8), Inches(0.45), fill_color=RGBColor(30, 41, 59), border_color=DANGER, border_width=1)
tf2 = decline.text_frame
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
run2 = tf2.paragraphs[0].add_run()
run2.text = "✗ Cancel"
run2.font.size = Pt(11)
run2.font.bold = True
run2.font.color.rgb = DANGER

add_footer_bar(slide)

# ==========================================
# SLIDE 8: BUSINESS IMPACT
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
cream_background(slide)

add_section_badge(slide, Inches(0.8), Inches(0.5), "BUSINESS IMPACT", SUCCESS)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Commercial Potential & Revenue Drivers",
             font_size=30, bold=True, color=CHARCOAL)

biz_cards = [
    ("💰", "Fee Income via Cross-Selling", 
     "Identifies idle balances and proactively recommends SBI Mutual Funds, Fixed Deposits, and Insurance. Increases digital product cross-sell conversion by contextualizing offers.",
     SBI_BLUE),
    ("📈", "Increased AUM", 
     "Goal-based micro-savings auto-routes round-ups and excess cash into SBI wealth products, growing total assets under management.",
     SUCCESS),
    ("❤️", "Customer Retention (LTV)", 
     "Proactive financial health monitoring (leakage detection, budget alerts) transforms SBI from a utility into a trusted wealth partner. Higher NPS, lower churn.",
     DANGER),
    ("🤖", "Support Cost Reduction", 
     "Autonomously resolves mandate blocks, transaction audits, budget checks, and transfers. Reduces call center ticket volume by 35-40%.",
     AMBER),
]

for i, (emoji, title, desc, color) in enumerate(biz_cards):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.15)
    y = Inches(2.0 + row * 2.6)
    
    card = add_shape(slide, x, y, Inches(5.9), Inches(2.3), fill_color=WHITE, border_color=RGBColor(235, 220, 211), border_width=1)
    add_shape(slide, x, y, Inches(5.9), Inches(0.06), fill_color=color)
    
    add_text_box(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 emoji, font_size=24)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.25), Inches(4.8), Inches(0.4),
                 title, font_size=16, bold=True, color=CHARCOAL)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.85), Inches(5.4), Inches(1.3),
                 desc, font_size=12, color=WARM_GRAY)

add_footer_bar(slide)

# ==========================================
# SLIDE 9: PROTOTYPE DEMO
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
cream_background(slide)

add_section_badge(slide, Inches(0.8), Inches(0.5), "LIVE PROTOTYPE", PURPLE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Fully Functional Interactive Dashboard — Ready for Phase 2",
             font_size=28, bold=True, color=CHARCOAL)

# 4 tab screenshots as placeholder cards
tabs = [
    ("📊 Overview Dashboard", "Real-time balance metrics, transaction log table, budget alert progress bars, and proactive agent notification cards."),
    ("🤖 AI Dost Assistant", "Conversational chat with live ReAct reasoning console showing THOUGHT → ACTION → OBSERVATION traces."),
    ("📋 Mandate Auditor", "Tabular audit of all recurring auto-debits with idle-day tracking, leakage risk badges, and one-click mandate blocking."),
    ("💹 Budgets & Goals", "Category budget trackers (Dining, Shopping, Rent), savings goal progress, and yield optimization recommendations."),
]

for i, (title, desc) in enumerate(tabs):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.15)
    y = Inches(2.0 + row * 2.5)
    
    card = add_shape(slide, x, y, Inches(5.9), Inches(2.2), fill_color=WHITE, border_color=RGBColor(235, 220, 211), border_width=1)
    
    add_text_box(slide, x + Inches(0.25), y + Inches(0.25), Inches(5.4), Inches(0.4),
                 title, font_size=16, bold=True, color=CHARCOAL)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.8), Inches(5.4), Inches(1.2),
                 desc, font_size=12, color=WARM_GRAY)

add_footer_bar(slide)

# ==========================================
# SLIDE 10: THANK YOU / CTA
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_background(slide)

# Decorative circle
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(2), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = DARK_BLUE
circle.line.fill.background()

add_text_box(slide, Inches(3), Inches(1.5), Inches(8), Inches(1),
             "Thank You", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(2.8), Inches(8), Inches(0.6),
             "SBI Dost — Autonomous Personal Finance Concierge",
             font_size=20, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.7), Inches(2.5), Inches(0.04), fill_color=SBI_BLUE)

add_text_box(slide, Inches(3), Inches(4.2), Inches(8), Inches(0.5),
             "Prototype Live  •  GitHub Repository Ready  •  3-Minute Demo Available",
             font_size=14, color=RGBColor(148, 163, 184), alignment=PP_ALIGN.CENTER)

# Contact info
add_text_box(slide, Inches(3), Inches(5.2), Inches(8), Inches(0.4),
             "Team: Jatin Garg", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3), Inches(5.7), Inches(8), Inches(0.4),
             "Pillar: Digital Engagement  •  Theme: Agentic AI & Emerging Tech",
             font_size=12, color=RGBColor(100, 116, 139), alignment=PP_ALIGN.CENTER)

add_footer_bar(slide)

# ==========================================
# SAVE
# ==========================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SBI_Dost_Idea_Deck.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
